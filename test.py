import os
from docx import Document
from pptx import Presentation
from fpdf import FPDF

# Ensure test_files directory exists
os.makedirs("test_files", exist_ok=True)

# Create an empty PDF
empty_pdf_path = "test_files/empty.pdf"
if not os.path.exists(empty_pdf_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.output(empty_pdf_path)
    print("Created:", empty_pdf_path)

# Create a PowerPoint file with no hyperlinks
no_links_ppt_path = "test_files/no_links.pptx"
if not os.path.exists(no_links_ppt_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(100, 100, 400, 200)
    textbox.text = "This slide has no hyperlinks."
    prs.save(no_links_ppt_path)
    print("Created:", no_links_ppt_path)

# Create a Word file with a complex table
complex_table_docx_path = "test_files/complex_table.docx"
if not os.path.exists(complex_table_docx_path):
    doc = Document()
    doc.add_paragraph("This document contains a complex table.")

    table = doc.add_table(rows=3, cols=3)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(0, 2).text = "Header 3"

    table.cell(1, 0).text = "Row 1, Col 1"
    table.cell(1, 1).merge(table.cell(2, 1))  # Merging cells to simulate complexity
    table.cell(1, 1).text = "Merged Cell"
    table.cell(1, 2).text = "Row 1, Col 3"
    table.cell(2, 0).text = "Row 2, Col 1"
    table.cell(2, 2).text = "Row 2, Col 3"

    doc.save(complex_table_docx_path)
    print("Created:", complex_table_docx_path)
