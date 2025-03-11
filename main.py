from abc import ABC, abstractmethod
from typing import List, Dict, Any
import os
import pdfplumber
import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from pptx import Presentation
import sqlite3
import csv
import json
import io

class FileLoader(ABC):
    """
    Abstract base class for loading and processing files.
    Provides methods for extracting text, links, images, and tables.
    """
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.metadata = self.extract_metadata()

    @abstractmethod
    def extract_text(self) -> Dict[str, Any]:
        pass

    @abstractmethod
    def extract_links(self) -> List[Dict[str, Any]]:
        pass

    @abstractmethod
    def extract_images(self) -> List[Dict[str, Any]]:
        pass

    @abstractmethod
    def extract_tables(self) -> List[Dict[str, Any]]:
        pass

    def extract_metadata(self) -> Dict[str, Any]:
    """
    Extracts metadata from the file.

    Returns:
        Dict[str, Any]: A dictionary containing:
            - file_name (str): Name of the file.
            - file_size (int): Size of the file in bytes.
            - file_type (str): File extension (e.g., .pdf, .docx, .pptx).

    Raises:
        FileNotFoundError: If the file does not exist.
    """
        return {
            "file_name": os.path.basename(self.file_path),
            "file_size": os.path.getsize(self.file_path),
            "file_type": os.path.splitext(self.file_path)[1]
        }

class PDFLoader(FileLoader):
    """
    Class for loading and extracting content from PDF files.
    """
    def extract_text(self) -> Dict[str, Any]:
    """
    Extracts text from a PDF file along with metadata.

    Returns:
        Dict[str, Any]: A dictionary containing:
            - text (dict[int, list[str]]): Text content indexed by page number.
            - metadata (dict): Metadata including:
                - font_styles (list[dict]): List of font styles used.
                - headings (dict[int, list[str]]): Headings per page.

    Notes:
        - Font size > 12 is assumed to be a heading.
        - Ignores empty text blocks.

    Example:
        >>> loader = PDFLoader("sample.pdf")
        >>> data = loader.extract_text()
        >>> print(data["text"][1])  # Prints text from page 1
    """
        text = {}
        headings = {}
        font_styles = []
        doc = fitz.open(self.file_path)
        
        for i, page in enumerate(doc):
            page_text = []
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        font_size = span["size"]
                        font_style = span["font"]
                        text_content = span["text"].strip()
                        if text_content:  # Ignore empty text
                            page_text.append(text_content)
                            font_styles.append({"page_number": i + 1, "text": text_content, "font": font_style, "size": font_size})
                            
                            if font_size > 12:  # Assuming headings have larger font sizes
                                if i + 1 not in headings:
                                    headings[i + 1] = []
                                headings[i + 1].append(text_content)
            text[i + 1] = page_text
        
        return {"text": text, "metadata": {"font_styles": font_styles, "headings": headings}}
    
    def extract_links(self) -> List[Dict[str, Any]]:
    """
    Extracts hyperlinks from the PDF file.

    Returns:
        List[Dict[str, Any]]: A list of dictionaries containing:
            - url (str): The hyperlink URL.
            - page_number (int): The page number where the link appears.

    Example:
        >>> loader = PDFLoader("sample.pdf")
        >>> links = loader.extract_links()
        >>> print(links)
        [{'url': 'https://example.com', 'page_number': 1}]

    """
        doc = fitz.open(self.file_path)
        links = []
        
        for i, page in enumerate(doc):
            page_links = page.get_links()
            
            for link in page_links:
                if 'uri' in link:
                    links.append({"url": link['uri'], "page_number": i + 1})
        
        return links
    
    def extract_images(self) -> List[Dict[str, Any]]:
    """
    Extracts images from the PDF and saves them as separate files.

    Returns:
        List[Dict[str, Any]]: A list of extracted images with:
            - page_number (int): Page where the image was found.
            - image_path (str): File path where the image was saved.

    Raises:
        IOError: If there is an issue saving the image.
    """

        doc = fitz.open(self.file_path)
        images = []
        output_dir = os.path.join("output", os.path.splitext(os.path.basename(self.file_path))[0], "images")
        os.makedirs(output_dir, exist_ok=True)

        for i, page in enumerate(doc):
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                img_ext = base_image["ext"]
                img_path = os.path.join(output_dir, f"page_{i+1}_img_{img_index}.{img_ext}")
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                images.append({"page_number": i + 1, "image_path": img_path})
        
        return images
    
    def extract_tables(self) -> List[Dict[str, Any]]:
    """
    Extracts tables from the PDF file and saves them as CSV.

    Returns:
        List[Dict[str, Any]]: A list of extracted tables with:
            - page_number (int): Page where the table was found.
            - table (list[list[str]]): Table data as a nested list.

    Notes:
        - Tables are saved as CSV files in the output directory.
    """
        tables = []
        output_dir = os.path.join("output", os.path.splitext(os.path.basename(self.file_path))[0])
        os.makedirs(output_dir, exist_ok=True)
        tables_path = os.path.join(output_dir, "extracted_tables.csv")
        
        with pdfplumber.open(self.file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                extracted_table = page.extract_table()
                if extracted_table:
                    tables.append({"page_number": i + 1, "table": extracted_table})
                    with open(tables_path, "w", newline="", encoding="utf-8") as f:
                        writer = csv.writer(f)
                        writer.writerows(extracted_table)
        return tables

class PPTLoader(FileLoader):
    """
    Class for loading and extracting content from PPT files.
    """
    def extract_text(self) -> Dict[str, Any]:
        """
Extracts text from a PowerPoint file.

Returns:
    Dict[str, Any]: Extracted text indexed by slide number.
"""

        text = {}
        headings = {}
        font_styles = []
        
        prs = Presentation(self.file_path)
        
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            slide_text = []
            
            # Extract title as heading
            if slide.shapes.title and hasattr(slide.shapes.title, "text_frame"):
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    slide_text.append(title_text)
                    
                    if slide_number not in headings:
                        headings[slide_number] = []
                    headings[slide_number].append(title_text)
                    
                    # Get better font information from title shape
                    for paragraph in slide.shapes.title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font_styles.append({
                                "page_number": slide_number,
                                "text": run.text.strip(),
                                "font": run.font.name if hasattr(run.font, "name") else "Default",
                                "size": run.font.size / 12700 if hasattr(run.font, "size") and run.font.size else 24  # Convert EMU to points
                            })
            
            # Extract text from other shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    # Handle each paragraph in text frames
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = paragraph.text.strip()
                        if paragraph_text and shape != slide.shapes.title:
                            slide_text.append(paragraph_text)
                            
                            # Get detailed font information
                            for run in paragraph.runs:
                                if run.text.strip():
                                    font_styles.append({
                                        "page_number": slide_number,
                                        "text": run.text.strip(),
                                        "font": run.font.name if hasattr(run.font, "name") else "Default",
                                        "size": run.font.size / 12700 if hasattr(run.font, "size") and run.font.size else 12  # Convert EMU to points
                                    })
            
            text[slide_number] = slide_text
        
        return {"text": text, "metadata": {"font_styles": font_styles, "headings": headings}}
    
    def extract_links(self) -> List[Dict[str, Any]]:
        """
Extracts hyperlinks from a PowerPoint file.

Returns:
    List[Dict[str, Any]]: List of extracted hyperlinks with:
        - url (str): The hyperlink URL.
        - page_number (int): Slide number where the hyperlink appears.
        - text (str): Text associated with the hyperlink (if available).
"""
        prs = Presentation(self.file_path)
        links = []
        
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            
            # Check all shapes for hyperlinks
            for shape in slide.shapes:
                # Check click action hyperlinks
                if hasattr(shape, "click_action") and shape.click_action:
                    if hasattr(shape.click_action, "hyperlink") and shape.click_action.hyperlink:
                        if hasattr(shape.click_action.hyperlink, "address") and shape.click_action.hyperlink.address:
                            shape_text = ""
                            if hasattr(shape, "text_frame") and shape.text_frame:
                                shape_text = shape.text_frame.text.strip()
                            
                            links.append({
                                "url": shape.click_action.hyperlink.address,
                                "page_number": slide_number,
                                "text": shape_text,
                                "shape_name": shape.name if hasattr(shape, "name") else "Unknown"
                            })
                
                # Check text runs for hyperlinks
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, "hyperlink") and run.hyperlink and hasattr(run.hyperlink, "address") and run.hyperlink.address:
                                links.append({
                                    "url": run.hyperlink.address,
                                    "page_number": slide_number,
                                    "text": run.text.strip(),
                                    "shape_name": shape.name if hasattr(shape, "name") else "Unknown"
                                })
        
        return links
    
    def extract_images(self) -> List[Dict[str, Any]]:
        """
Extracts images from PowerPoint slides.

Returns:
    List[Dict[str, Any]]: List of extracted images with:
        - page_number (int): Slide number where the image appears.
        - image_path (str): File path where the image was saved.
"""
        prs = Presentation(self.file_path)
        images = []
        output_dir = os.path.join("output", os.path.splitext(os.path.basename(self.file_path))[0], "images")
        os.makedirs(output_dir, exist_ok=True)
        
        image_index = 0
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            
            for shape in slide.shapes:
                # Check for pictures (shape type 13)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                            image_stream = io.BytesIO(shape.image.blob)
                            image = Image.open(image_stream)
                            img_path = os.path.join(output_dir, f"slide_{slide_number}_img_{image_index}.png")
                            image.save(img_path)
                            
                            images.append({
                                "page_number": slide_number,
                                "image_path": img_path,
                                "alt_text": shape.alt_text if hasattr(shape, "alt_text") else ""
                            })
                            
                            image_index += 1
                    except Exception as e:
                        print(f"Error extracting image: {e}")
        
        return images
    
    def extract_tables(self) -> List[Dict[str, Any]]:
        """
Extracts tables from PowerPoint slides.

Returns:
    List[Dict[str, Any]]: List of extracted tables with:
        - page_number (int): Slide number where the table appears.
        - table (list[list[str]]): Table data as a nested list.
"""
        prs = Presentation(self.file_path)
        tables = []
        output_dir = os.path.join("output", os.path.splitext(os.path.basename(self.file_path))[0])
        os.makedirs(output_dir, exist_ok=True)
        
        table_index = 0
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    extracted_table = []
                    
                    for row in shape.table.rows:
                        table_row = []
                        for cell in row.cells:
                            if hasattr(cell, "text_frame") and cell.text_frame:
                                table_row.append(cell.text_frame.text.strip())
                            else:
                                table_row.append("")
                        extracted_table.append(table_row)
                    
                    # Only save non-empty tables
                    if any(any(cell for cell in row) for row in extracted_table):
                        tables.append({
                            "page_number": slide_number,
                            "table": extracted_table,
                            "table_index": table_index
                        })
                        
                        # Save extracted table to CSV
                        tables_path = os.path.join(output_dir, f"extracted_table_{table_index}.csv")
                        with open(tables_path, "w", newline="", encoding="utf-8") as f:
                            writer = csv.writer(f)
                            writer.writerows(extracted_table)
                        
                        table_index += 1
        
        return tables

class DOCXLoader(FileLoader):
    """
    Class for loading and extracting content from DOCX files.
    """
    def extract_text(self) -> Dict[str, Any]:
        """
Extracts text from a Word document.

Returns:
    Dict[str, Any]: A dictionary containing:
        - text (dict[int, list[str]]): Extracted text content.
        - metadata (dict): Metadata including:
            - font_styles (list[dict]): List of font styles used.
            - headings (dict[int, list[str]]): Headings per page.
"""
        text = {}
        headings = {}
        font_styles = []
        
        doc = Document(self.file_path)
        
        # Process paragraphs
        page_number = 1  # Since DOCX doesn't have explicit page numbers, we'll treat the whole document as one page
        page_text = []
        
        for para in doc.paragraphs:
            text_content = para.text.strip()
            if text_content:
                page_text.append(text_content)
                
                # Extract style information
                style_name = para.style.name if para.style else "Normal"
                font_size = 12  # Default size, can't easily get actual size
                
                # Collect font information
                font_styles.append({
                    "page_number": page_number,
                    "text": text_content,
                    "font": style_name,
                    "size": font_size
                })
                
                # Check if this is a heading
                if style_name.startswith('Heading'):
                    if page_number not in headings:
                        headings[page_number] = []
                    headings[page_number].append(text_content)
        
        text[page_number] = page_text
        
        return {"text": text, "metadata": {"font_styles": font_styles, "headings": headings}}
    
    def extract_links(self) -> List[Dict[str, Any]]:
        """
Extracts hyperlinks from a Word document.

Returns:
    List[Dict[str, Any]]: List of extracted hyperlinks with:
        - url (str): The hyperlink URL.
        - page_number (int): Page number where the hyperlink appears.
        - text (str): Text associated with the hyperlink (if available).
"""

        doc = Document(self.file_path)
        links = []
        page_number = 1  # Treat whole document as one page
        
        for para in doc.paragraphs:
            for run in para.runs:
                # Access the XML element directly
                if hasattr(run._element, 'xpath'):
                    # Find hyperlink references in the run
                    hyperlink_refs = run._element.xpath('.//w:hyperlink')
                    if hyperlink_refs:
                        for hyperlink in hyperlink_refs:
                            # Get relationship ID
                            rel_id = hyperlink.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                            if rel_id and rel_id in doc.part.rels:
                                target_url = doc.part.rels[rel_id].target_ref
                                # Get text from the hyperlink
                                hyperlink_text = ""
                                for text_element in hyperlink.xpath('.//w:t'):
                                    if text_element.text:
                                        hyperlink_text += text_element.text
                                
                                links.append({
                                    "url": target_url,
                                    "page_number": page_number,
                                    "text": hyperlink_text
                                })
        
        # Alternative method: parse directly from relationships
        for rel in doc.part.rels.values():
            if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink':
                if rel.target_ref.startswith('http'):
                    # Only add if not already found
                    if not any(link['url'] == rel.target_ref for link in links):
                        links.append({
                            "url": rel.target_ref,
                            "page_number": page_number,
                            "text": "Link"  # Default text since we can't match it to content
                        })
        
        return links
    
    def extract_images(self) -> List[Dict[str, Any]]:
        """
Extracts images from a Word document.

Returns:
    List[Dict[str, Any]]: List of extracted images with:
        - page_number (int): Page number where the image appears.
        - image_path (str): File path where the image was saved.
"""

        doc = Document(self.file_path)
        images = []
        output_dir = os.path.join("output", os.path.splitext(os.path.basename(self.file_path))[0], "images")
        os.makedirs(output_dir, exist_ok=True)
        
        image_index = 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    img_ext = rel.target_ref.split(".")[-1]
                    img_path = os.path.join(output_dir, f"image_{image_index}.{img_ext}")
                    
                    with open(img_path, "wb") as f:
                        f.write(image_data)
                    
                    images.append({
                        "page_number": 1,  # Assuming single page for simplicity
                        "image_path": img_path
                    })
                    
                    image_index += 1
                except Exception as e:
                    print(f"Error extracting image: {e}")
        
        return images
    
    def extract_tables(self) -> List[Dict[str, Any]]:
        """
Extracts tables from a Word document.

Returns:
    List[Dict[str, Any]]: List of extracted tables with:
        - page_number (int): Page number where the table appears.
        - table (list[list[str]]): Table data as a nested list.
"""

        doc = Document(self.file_path)
        tables = []
        output_dir = os.path.join("output", os.path.splitext(os.path.basename(self.file_path))[0])
        os.makedirs(output_dir, exist_ok=True)
        tables_path = os.path.join(output_dir, "extracted_tables.csv")
        
        for table_index, table in enumerate(doc.tables):
            extracted_table = []
            
            # Extract header row
            header_row = []
            for cell in table.rows[0].cells:
                header_row.append(cell.text.strip())
            extracted_table.append(header_row)
            
            # Extract data rows
            for row in table.rows[1:]:
                data_row = []
                for cell in row.cells:
                    data_row.append(cell.text.strip())
                extracted_table.append(data_row)
            
            tables.append({
                "page_number": 1,  # Assuming single page for simplicity
                "table": extracted_table,
                "table_index": table_index
            })
            
            # Save extracted table to CSV
            with open(tables_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                writer.writerows(extracted_table)
        
        return tables

class DataExtractor:
    """
    Wrapper class for extracting content using a FileLoader instance.
    """
    def __init__(self, file_loader: FileLoader):
        self.file_loader = file_loader
    
    def extract_text(self) -> Dict[str, Any]:
        return self.file_loader.extract_text()
    
    def extract_links(self) -> List[Dict[str, Any]]:
        return self.file_loader.extract_links()
    
    def extract_images(self) -> List[Dict[str, Any]]:
        return self.file_loader.extract_images()
    
    def extract_tables(self) -> List[Dict[str, Any]]:
        return self.file_loader.extract_tables()

class Storage(ABC):
    @abstractmethod
    def save(self, data: Dict[str, Any], file_name: str):
        pass

class FileStorage(Storage):
"""
Saves extracted document data to local files.

Args:
    data (Dict[str, Any]): Extracted text, links, images, and tables.
    file_name (str): Name of the file where data will be saved.

Notes:
    - Text is saved as a `.txt` file.
    - Links are saved as `.csv`.
    - Font styles are saved separately.
"""

    def save(self, data: Dict[str, Any], file_name: str):
        output_dir = os.path.join("output", file_name)
        os.makedirs(output_dir, exist_ok=True)

        text_path = os.path.join(output_dir, "extracted_text.txt")
        with open(text_path, "w", encoding="utf-8") as f:
            for page, content in data.get("text", {}).get("text", {}).items():
                f.write(f"Page {page}\n")
                f.write("\n".join(content) + "\n\n")
        
        headings_path = os.path.join(output_dir, "headings.txt")
        with open(headings_path, "w", encoding="utf-8") as f:
            for page, headings in data.get("text", {}).get("metadata", {}).get("headings", {}).items():
                f.write(f"Page {page}\n")
                f.write("\n".join(headings) + "\n\n")
        
        links_path = os.path.join(output_dir, "extracted_links.csv")
        with open(links_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Page Number", "URL"])
            for link in data.get("links", []):
                writer.writerow([link["page_number"], link["url"]])
        
        font_styles_path = os.path.join(output_dir, "font_styles.csv")
        with open(font_styles_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Page Number", "Text", "Font", "Size"])
            for font in data.get("text", {}).get("metadata", {}).get("font_styles", []):
                if font["text"].strip():  # Ignore empty text
                    writer.writerow([font["page_number"], font["text"], font["font"], font["size"]])

class SQLStorage(Storage):
"""
Saves extracted document data into an SQLite database.

Args:
    data (Dict[str, Any]): Extracted text, links, images, and tables.
    file_name (str): Name of the document.

Raises:
    sqlite3.DatabaseError: If an error occurs during the database operation.
"""
    def __init__(self, db_path="document_data.db"):
        """
        Initialize SQLStorage with database connection.
        
        Args:
            db_path (str): Path to the SQLite database file
        """
        self.db_path = db_path
        self._create_tables()
    
    def _create_tables(self):
        """Create the necessary database tables if they don't exist."""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # Create tables for documents and extracted data
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_name TEXT NOT NULL,
            file_path TEXT NOT NULL,
            file_size INTEGER,
            file_type TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
        ''')
        
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_text (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER,
            page_number INTEGER,
            content TEXT,
            FOREIGN KEY (document_id) REFERENCES documents (id)
        )
        ''')
        
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_headings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER,
            page_number INTEGER,
            heading TEXT,
            FOREIGN KEY (document_id) REFERENCES documents (id)
        )
        ''')
        
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_links (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER,
            page_number INTEGER,
            url TEXT,
            link_text TEXT,
            FOREIGN KEY (document_id) REFERENCES documents (id)
        )
        ''')
        
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_images (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER,
            page_number INTEGER,
            image_path TEXT,
            alt_text TEXT,
            FOREIGN KEY (document_id) REFERENCES documents (id)
        )
        ''')
        
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_tables (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER,
            page_number INTEGER,
            table_index INTEGER,
            table_data TEXT,  
            FOREIGN KEY (document_id) REFERENCES documents (id)
        )
        ''')
        
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_font_styles (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER,
            page_number INTEGER,
            text TEXT,
            font TEXT,
            size REAL,
            FOREIGN KEY (document_id) REFERENCES documents (id)
        )
        ''')
        
        conn.commit()
        conn.close()
    
    def save(self, data: Dict[str, Any], file_name: str):
        """
        Save extracted document data to the SQLite database.
        
        Args:
            data (Dict[str, Any]): Dictionary containing extracted document data
            file_name (str): Name of the file
        """
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        try:
            # Insert document metadata
            cursor.execute(
                "INSERT INTO documents (file_name, file_path, file_size, file_type) VALUES (?, ?, ?, ?)",
                (
                    file_name,
                    data.get("text", {}).get("metadata", {}).get("file_path", ""),
                    data.get("text", {}).get("metadata", {}).get("file_size", 0),
                    data.get("text", {}).get("metadata", {}).get("file_type", "")
                )
            )
            document_id = cursor.lastrowid
            
            # Insert document text
            for page_num, text_content in data.get("text", {}).get("text", {}).items():
                # Join all text elements into a single string for storage
                page_text = "\n".join(text_content) if isinstance(text_content, list) else str(text_content)
                cursor.execute(
                    "INSERT INTO document_text (document_id, page_number, content) VALUES (?, ?, ?)",
                    (document_id, page_num, page_text)
                )
            
            # Insert headings
            for page_num, headings in data.get("text", {}).get("metadata", {}).get("headings", {}).items():
                for heading in headings:
                    cursor.execute(
                        "INSERT INTO document_headings (document_id, page_number, heading) VALUES (?, ?, ?)",
                        (document_id, page_num, heading)
                    )
            
            # Insert links
            for link in data.get("links", []):
                cursor.execute(
                    "INSERT INTO document_links (document_id, page_number, url, link_text) VALUES (?, ?, ?, ?)",
                    (
                        document_id, 
                        link.get("page_number", 0), 
                        link.get("url", ""), 
                        link.get("text", "")
                    )
                )
            
            # Insert images
            for image in data.get("images", []):
                cursor.execute(
                    "INSERT INTO document_images (document_id, page_number, image_path, alt_text) VALUES (?, ?, ?, ?)",
                    (
                        document_id, 
                        image.get("page_number", 0), 
                        image.get("image_path", ""), 
                        image.get("alt_text", "")
                    )
                )
            
            # Insert tables (converting table data to JSON string)
            for table in data.get("tables", []):
                table_data = json.dumps(table.get("table", []))
                cursor.execute(
                    "INSERT INTO document_tables (document_id, page_number, table_index, table_data) VALUES (?, ?, ?, ?)",
                    (
                        document_id, 
                        table.get("page_number", 0), 
                        table.get("table_index", 0), 
                        table_data
                    )
                )
            
            # Insert font styles
            for font_style in data.get("text", {}).get("metadata", {}).get("font_styles", []):
                if font_style.get("text", "").strip():  # Ignore empty text
                    cursor.execute(
                        "INSERT INTO document_font_styles (document_id, page_number, text, font, size) VALUES (?, ?, ?, ?, ?)",
                        (
                            document_id, 
                            font_style.get("page_number", 0), 
                            font_style.get("text", ""), 
                            font_style.get("font", ""), 
                            font_style.get("size", 0)
                        )
                    )
            
            conn.commit()
            print(f"Successfully saved data for '{file_name}' to database")
            
        except Exception as e:
            conn.rollback()
            print(f"Error saving data to database: {e}")
        finally:
            conn.close()
    
    def query_document(self, file_name=None, document_id=None):
        """
        Query document data from the database.
        
        Args:
            file_name (str, optional): Name of the file to query
            document_id (int, optional): ID of the document to query
            
        Returns:
            Dict[str, Any]: Dictionary containing document data
        """
        conn = sqlite3.connect(self.db_path)
        conn.row_factory = sqlite3.Row  # Return rows as dictionaries
        cursor = conn.cursor()
        
        result = {}
        
        try:
            # Query document by file_name or document_id
            if file_name:
                cursor.execute("SELECT * FROM documents WHERE file_name = ?", (file_name,))
            elif document_id:
                cursor.execute("SELECT * FROM documents WHERE id = ?", (document_id,))
            else:
                return {"error": "Either file_name or document_id must be provided"}
            
            document = cursor.fetchone()
            if not document:
                return {"error": "Document not found"}
            
            document_id = document["id"]
            result["document"] = dict(document)
            
            # Query document text
            cursor.execute("SELECT * FROM document_text WHERE document_id = ? ORDER BY page_number", (document_id,))
            result["text"] = [dict(row) for row in cursor.fetchall()]
            
            # Query headings
            cursor.execute("SELECT * FROM document_headings WHERE document_id = ? ORDER BY page_number", (document_id,))
            result["headings"] = [dict(row) for row in cursor.fetchall()]
            
            # Query links
            cursor.execute("SELECT * FROM document_links WHERE document_id = ? ORDER BY page_number", (document_id,))
            result["links"] = [dict(row) for row in cursor.fetchall()]
            
            # Query images
            cursor.execute("SELECT * FROM document_images WHERE document_id = ? ORDER BY page_number", (document_id,))
            result["images"] = [dict(row) for row in cursor.fetchall()]
            
            # Query tables
            cursor.execute("SELECT * FROM document_tables WHERE document_id = ? ORDER BY page_number, table_index", (document_id,))
            tables = []
            for row in cursor.fetchall():
                row_dict = dict(row)
                # Parse JSON table data
                try:
                    row_dict["table"] = json.loads(row_dict["table_data"])
                    del row_dict["table_data"]
                except:
                    row_dict["table"] = []
                tables.append(row_dict)
            result["tables"] = tables
            
            # Query font styles
            cursor.execute("SELECT * FROM document_font_styles WHERE document_id = ? ORDER BY page_number", (document_id,))
            result["font_styles"] = [dict(row) for row in cursor.fetchall()]
            
        except Exception as e:
            print(f"Error querying database: {e}")
            result["error"] = str(e)
        finally:
            conn.close()
        
        return result
    
    def list_documents(self):
        """
        List all documents in the database.
        
        Returns:
            List[Dict[str, Any]]: List of document metadata
        """
        conn = sqlite3.connect(self.db_path)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        
        try:
            cursor.execute("SELECT * FROM documents ORDER BY created_at DESC")
            documents = [dict(row) for row in cursor.fetchall()]
            return documents
        except Exception as e:
            print(f"Error listing documents: {e}")
            return []
        finally:
            conn.close()
    
    def delete_document(self, document_id):
        """
        Delete a document and all its associated data from the database.
        
        Args:
            document_id (int): ID of the document to delete
            
        Returns:
            bool: True if successful, False otherwise
        """
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        try:
            # Delete associated data first (foreign key constraints)
            for table in ["document_text", "document_headings", "document_links", 
                          "document_images", "document_tables", "document_font_styles"]:
                cursor.execute(f"DELETE FROM {table} WHERE document_id = ?", (document_id,))
            
            # Delete the document
            cursor.execute("DELETE FROM documents WHERE id = ?", (document_id,))
            conn.commit()
            return True
        except Exception as e:
            conn.rollback()
            print(f"Error deleting document: {e}")
            return False
        finally:
            conn.close()

#Sample usage
# pptx_loader = PPTLoader("input/Chatfolio.pptx")
# extractor = DataExtractor(pptx_loader)

# file_name = os.path.splitext(os.path.basename(pptx_loader.file_path))[0]
# data = {
#     "text": extractor.extract_text(),
#     "links": extractor.extract_links(),
#     "images": extractor.extract_images(),
#     "tables": extractor.extract_tables()
# }

# storage = FileStorage()
# storage.save(data, file_name)
# storage = SQLStorage("documents.db")  
# storage.save(data, file_name)
